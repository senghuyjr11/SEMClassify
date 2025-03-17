import os
from pptx import Presentation
from PIL import Image
from io import BytesIO
import re

# Define PowerPoint file path
pptx_file = "dataset/RealCorrosionData.pptx"
output_root = "dataset"

# Load PowerPoint presentation
presentation = Presentation(pptx_file)

# Extract images from slides
for slide_number, slide in enumerate(presentation.slides, start=1):
    slide_text = " ".join(
        [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.has_text_frame]).lower()

    # Extract category from text using regex
    match = re.search(r"(an_\d+) after (\d+\s*(?:hr|days))", slide_text)
    if match:
        annealing_time, exposure_duration = match.groups()
        # Standardize folder names
        exposure_duration = exposure_duration.replace(" ", "")  # Convert '7 days' to '7days'
    else:
        annealing_time, exposure_duration = "Uncategorized", "Unknown"
        print(f"Slide {slide_number} text not recognized. Saving to Uncategorized folder.")

    output_folder = os.path.join(output_root, annealing_time, exposure_duration)
    os.makedirs(output_folder, exist_ok=True)

    image_count = 0
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            image = shape.image
            image_bytes = BytesIO(image.blob)
            img = Image.open(image_bytes)

            # Determine image format
            img_format = img.format.lower()
            if img_format not in ["jpeg", "png", "gif", "bmp"]:
                img_format = "png"  # Default to PNG if format is unknown

            # Save image
            img_filename = os.path.join(output_folder, f"Img_Slide{slide_number}_No{image_count}.{img_format}")
            img.save(img_filename)
            print(f"Saved: {img_filename}")
            image_count += 1

print("Extraction complete. Check the dataset folder.")


# Verify dataset completeness
def verify_dataset(dataset_root, expected_structure):
    missing_folders = []
    missing_images = []

    for annealing, durations in expected_structure.items():
        for duration in durations:
            folder_path = os.path.join(dataset_root, annealing, duration)
            if not os.path.exists(folder_path):
                missing_folders.append(folder_path)
            else:
                images = os.listdir(folder_path)
                if len(images) == 0:
                    missing_images.append(folder_path)

    if missing_folders:
        print("Missing folders:")
        for folder in missing_folders:
            print(f" - {folder}")
    else:
        print("All expected folders are present.")

    if missing_images:
        print("Folders with no images:")
        for folder in missing_images:
            print(f" - {folder}")
    else:
        print("All folders contain images.")


# Define expected dataset structure
expected_structure = {
    "AN_30": ["1hr", "7days", "14days", "21days"],
    "AN_60": ["1hr", "7days", "14days", "21days"],
    "AN_180": ["1hr", "7days", "14days", "21days"],
}

# Run verification
verify_dataset(output_root, expected_structure)
